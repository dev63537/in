import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

try:
    prs = Presentation()
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 1. Background
    background_img = r"C:\Users\meetc\.gemini\antigravity\brain\ce5fc048-f9a1-4612-accb-65679e42ee24\cyber_background_1783954260555.png"
    slide.shapes.add_picture(background_img, 0, 0, prs.slide_width, prs.slide_height)

    # 2. Top Header
    # Top Left
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "InventMitra"
    p.font.name = 'Montserrat'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 255, 255) # Cyan

    # Top Right
    txBox_tr = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5))
    tf_tr = txBox_tr.text_frame
    p_tr = tf_tr.paragraphs[0]
    p_tr.text = "01 / INVENTMITRA"
    p_tr.font.name = 'Inter'
    p_tr.font.size = Pt(12)
    p_tr.font.color.rgb = RGBColor(200, 200, 200)
    p_tr.alignment = PP_ALIGN.RIGHT

    # 3. Left Side Content (40%)
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5), Inches(1.5))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "InventMitra"
    p_title.font.name = 'Poppins'
    p_title.font.size = Pt(60)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)

    # Cyan Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(2.6), Inches(3.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 255, 255)
    line.line.fill.background()

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5), Inches(0.5))
    tf_tag = tagline_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "Your Trusted Inventory Partner"
    p_tag.font.name = 'Inter'
    p_tag.font.size = Pt(20)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(0, 200, 255)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(5), Inches(0.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Smart Inventory Management System"
    p_sub.font.name = 'Montserrat'
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = RGBColor(220, 220, 220)

    # Technology Pills
    techs = ["HTML", "CSS", "JavaScript", "PHP", "MySQL", "Bootstrap"]
    start_x = 0.85
    start_y = 3.9
    pill_width = 0.95
    pill_height = 0.35
    for tech in techs:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(start_y), Inches(pill_width), Inches(pill_height))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(10, 30, 60) # Dark navy
        pill.line.color.rgb = RGBColor(0, 255, 255) # Cyan glow border
        
        tf_pill = pill.text_frame
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = tech
        p_pill.font.name = 'Inter'
        p_pill.font.size = Pt(10)
        p_pill.font.color.rgb = RGBColor(255, 255, 255)
        p_pill.alignment = PP_ALIGN.CENTER
        start_x += (pill_width + 0.1)

    # Tiny Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(4.5), Inches(1))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A modern web platform designed to simplify inventory management through secure digital records, product tracking, supplier management, purchase monitoring and reporting."
    p_desc.font.name = 'Inter'
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = RGBColor(180, 180, 180)

    # Badges
    badge_y = 5.5
    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(badge_y), Inches(2.2), Inches(0.4))
    b1.fill.solid()
    b1.fill.fore_color.rgb = RGBColor(0, 120, 255) # Electric Blue
    b1.line.fill.background()
    tf_b1 = b1.text_frame
    p_b1 = tf_b1.paragraphs[0]
    p_b1.text = "BCA Semester 5 Project"
    p_b1.font.name = 'Montserrat'
    p_b1.font.size = Pt(10)
    p_b1.font.bold = True
    p_b1.font.color.rgb = RGBColor(255, 255, 255)
    p_b1.alignment = PP_ALIGN.CENTER

    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(badge_y), Inches(1.2), Inches(0.4))
    b2.fill.solid()
    b2.fill.fore_color.rgb = RGBColor(0, 200, 100) # Greenish/Neon
    b2.line.fill.background()
    tf_b2 = b2.text_frame
    p_b2 = tf_b2.paragraphs[0]
    p_b2.text = "Viva Ready"
    p_b2.font.name = 'Montserrat'
    p_b2.font.size = Pt(10)
    p_b2.font.bold = True
    p_b2.font.color.rgb = RGBColor(255, 255, 255)
    p_b2.alignment = PP_ALIGN.CENTER

    # 4. Right Side (Dashboard & Illustration)
    dashboard_img = r"C:\Users\meetc\.gemini\antigravity\brain\ce5fc048-f9a1-4612-accb-65679e42ee24\dashboard_ui_1783954279356.png"
    iso_img = r"C:\Users\meetc\.gemini\antigravity\brain\ce5fc048-f9a1-4612-accb-65679e42ee24\isometric_warehouse_1783954236119.png"

    # Place Dashboard Image (Top Right)
    slide.shapes.add_picture(dashboard_img, Inches(5.8), Inches(1.0), width=Inches(7.0))

    # Place Isometric Warehouse Image (Bottom Right)
    slide.shapes.add_picture(iso_img, Inches(6.5), Inches(4.0), width=Inches(5.5))

    # 5. Footer (Glass Navigation Bar)
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.4))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(10, 15, 30) # Very dark
    footer_bar.line.color.rgb = RGBColor(0, 255, 255)

    def add_footer_text(text, x, width, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(6.8), Inches(width), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Inter'
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = align

    add_footer_text("BCA Semester 5 Project", 0.6, 2.5, PP_ALIGN.LEFT)
    add_footer_text("Presented by: Meet Chavda", 3.5, 2.5, PP_ALIGN.CENTER)
    add_footer_text("Guide: Vipulkumar Baldha", 6.5, 2.5, PP_ALIGN.CENTER)
    add_footer_text("Kamani Science & Prataprai arts college", 9.5, 3.2, PP_ALIGN.RIGHT)

    # Save
    prs.save(r"C:\Users\meetc\Videos\project presention\InventMitra_Cover.pptx")
    print("SUCCESS")
except Exception as e:
    print(f"ERROR: {e}")
